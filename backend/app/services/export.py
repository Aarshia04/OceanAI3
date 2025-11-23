from typing import Iterable

from docx import Document
from pptx import Presentation

from app.models import GeneratedContent, Project


class DocExporter:
    def export(self, project: Project, content: Iterable[GeneratedContent], output_path: str, format: str = "docx") -> str:
        if format == "docx":
            self._export_docx(project, content, output_path)
        else:
            self._export_pptx(project, content, output_path)
        return output_path

    def _export_docx(self, project: Project, content: Iterable[GeneratedContent], output_path: str) -> None:
        doc = Document()
        doc.add_heading(project.name, 0)
        for entry in content:
            doc.add_heading(entry.section, level=1)
            doc.add_paragraph(entry.content)
        doc.save(output_path)

    def _export_pptx(self, project: Project, content: Iterable[GeneratedContent], output_path: str) -> None:
        presentation = Presentation()
        title_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = project.name
        for entry in content:
            layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = entry.section
            slide.placeholders[1].text = entry.content
        presentation.save(output_path)
