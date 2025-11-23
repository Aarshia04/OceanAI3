from dataclasses import dataclass, field
from io import BytesIO
from typing import List

from docx import Document
from pptx import Presentation
from pptx.util import Pt


@dataclass
class DocumentSection:
    title: str
    paragraphs: List[str] = field(default_factory=list)


def export_docx(sections: list[DocumentSection]) -> BytesIO:
    """Build a DOCX document from the provided sections.

    Args:
        sections: The sections to include in the document.

    Returns:
        A ``BytesIO`` object containing the DOCX bytes.
    """
    document = Document()

    for section in sections:
        document.add_heading(section.title, level=1)
        for paragraph in section.paragraphs:
            document.add_paragraph(paragraph)

    output = BytesIO()
    document.save(output)
    output.seek(0)
    return output


def export_pptx(sections: list[DocumentSection]) -> BytesIO:
    """Build a PPTX slide deck from the provided sections.

    Args:
        sections: The sections to include in the deck.

    Returns:
        A ``BytesIO`` object containing the PPTX bytes.
    """
    presentation = Presentation()

    title_layout = presentation.slide_layouts[1]

    for section in sections:
        slide = presentation.slides.add_slide(title_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title_placeholder.text = section.title
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        for paragraph in section.paragraphs:
            p = text_frame.add_paragraph()
            p.text = paragraph
            p.font.size = Pt(18)
            p.level = 0

        if not section.paragraphs:
            text_frame.text = ""

    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    return output
