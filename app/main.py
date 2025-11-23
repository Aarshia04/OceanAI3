from datetime import timedelta
from io import BytesIO

from fastapi import Depends, FastAPI, HTTPException, status
from fastapi.responses import StreamingResponse
from fastapi.security import OAuth2PasswordRequestForm
from docx import Document
from pptx import Presentation

from app import llm
from app.auth import (
    ACCESS_TOKEN_EXPIRE_MINUTES,
    create_access_token,
    get_current_user,
    get_password_hash,
    verify_password,
)
from app.models import Project, Refinement, User
from app.schemas import (
    ContentGenerateRequest,
    ProjectCreateRequest,
    ProjectResponse,
    RefinementHistoryResponse,
    RefinementRequest,
    RegisterRequest,
    SectionContentResponse,
    TokenResponse,
    UpdateDocsRequest,
)
from app.storage import store

app = FastAPI(title="OceanAI Content API")


def authenticate_user(username: str, password: str) -> User | None:
    user = store.get_user(username)
    if not user:
        return None
    if not verify_password(password, user.hashed_password):
        return None
    return user


@app.post("/register", response_model=TokenResponse, status_code=status.HTTP_201_CREATED)
def register(payload: RegisterRequest):
    if store.get_user(payload.username):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="User already exists")

    hashed_password = get_password_hash(payload.password)
    store.create_user(payload.username, hashed_password)

    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    token = create_access_token({"sub": payload.username}, access_token_expires)
    return TokenResponse(access_token=token)


@app.post("/login", response_model=TokenResponse)
def login(form_data: OAuth2PasswordRequestForm = Depends()):
    user = authenticate_user(form_data.username, form_data.password)
    if not user:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Invalid credentials")

    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    token = create_access_token({"sub": user.username}, access_token_expires)
    return TokenResponse(access_token=token)


@app.post("/projects", response_model=ProjectResponse, status_code=status.HTTP_201_CREATED)
def create_project(payload: ProjectCreateRequest, _: User = Depends(get_current_user)):
    project = store.create_project(payload.name)
    return ProjectResponse(id=project.id, name=project.name, sections=project.sections)


@app.put("/projects/{project_id}/docs", response_model=ProjectResponse)
def update_project_docs(project_id: str, payload: UpdateDocsRequest, _: User = Depends(get_current_user)):
    project = _get_project_or_404(project_id)
    project.sections = payload.sections
    for section in payload.sections:
        project.section_content.setdefault(section, "")
    store.save_project(project)
    return ProjectResponse(id=project.id, name=project.name, sections=project.sections)


@app.post("/projects/{project_id}/sections/{section}/generate", response_model=SectionContentResponse)
def generate_content(project_id: str, section: str, payload: ContentGenerateRequest, _: User = Depends(get_current_user)):
    project = _get_project_or_404(project_id)
    if section not in project.sections:
        project.sections.append(section)
    generated = llm.generate_section_content(section, payload.prompt)
    project.section_content[section] = generated
    store.save_project(project)
    return SectionContentResponse(project_id=project.id, section=section, content=generated)


@app.post("/projects/{project_id}/sections/{section}/refine", response_model=SectionContentResponse)
def refine_content(project_id: str, section: str, payload: RefinementRequest, _: User = Depends(get_current_user)):
    project = _get_project_or_404(project_id)
    existing = project.section_content.get(section)
    if existing is None:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Section content not found")

    base_content = payload.content or existing
    refined_content = llm.refine_section_content(base_content, payload.feedback)
    refinement = Refinement(
        section=section,
        content_before=base_content,
        content_after=refined_content,
        feedback=payload.feedback,
    )
    project.section_content[section] = refined_content
    project.refinement_history.append(refinement)
    store.save_project(project)
    return SectionContentResponse(project_id=project.id, section=section, content=refined_content, refined=True)


@app.get("/projects/{project_id}/refinements", response_model=RefinementHistoryResponse)
def get_refinement_history(project_id: str, _: User = Depends(get_current_user)):
    project = _get_project_or_404(project_id)
    history = [
        {
            "section": record.section,
            "content_before": record.content_before,
            "content_after": record.content_after,
            "feedback": record.feedback,
            "refined_at": record.refined_at,
        }
        for record in project.refinement_history
    ]
    return RefinementHistoryResponse(project_id=project.id, history=history)


@app.get("/projects/{project_id}/export/docx")
def export_docx(project_id: str, _: User = Depends(get_current_user)):
    project = _get_project_or_404(project_id)
    doc = Document()
    doc.add_heading(project.name, 0)
    for section in project.sections:
        doc.add_heading(section, level=1)
        doc.add_paragraph(project.section_content.get(section, ""))
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return StreamingResponse(buffer, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", headers={"Content-Disposition": f"attachment; filename={project.name}.docx"})


@app.get("/projects/{project_id}/export/pptx")
def export_pptx(project_id: str, _: User = Depends(get_current_user)):
    project = _get_project_or_404(project_id)
    presentation = Presentation()
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = project.name
    slide.placeholders[1].text = "Project overview"

    bullet_layout = presentation.slide_layouts[1]
    for section in project.sections:
        slide = presentation.slides.add_slide(bullet_layout)
        slide.shapes.title.text = section
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = project.section_content.get(section, "")

    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return StreamingResponse(buffer, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": f"attachment; filename={project.name}.pptx"})


def _get_project_or_404(project_id: str) -> Project:
    project = store.get_project(project_id)
    if not project:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Project not found")
    return project
